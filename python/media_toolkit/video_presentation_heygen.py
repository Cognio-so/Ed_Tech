import os
import time
import json
import logging
import mimetypes
from typing import List, Dict, Optional

import requests
from dotenv import load_dotenv
from pptx import Presentation

# --- OpenAI SDK Setup ---
try:
    from openai import OpenAI
except ImportError:
    raise ImportError("OpenAI SDK not found. Install with: pip install openai")

load_dotenv()

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")


class PPTXToHeyGenVideo:
    def __init__(
        self,
        heygen_api_key: Optional[str] = None,
        openai_api_key: Optional[str] = None,
        avatar_id: Optional[str] = None,
        voice_id: Optional[str] = None,
        model: str = "gpt-4.1",
        width: int = 1280,
        height: int = 720,
        background_color: str = "#FFFFFF",
        request_timeout_s: int = 70,
        poll_interval_s: int = 5,
        pptx_avatar_id: Optional[str] = None,
        pptx_voice_id: Optional[str] = None,
        # --- MODIFICATION: Hardcoded to always use slides as background ---
        use_slides_as_background: bool = True,
    ):
        self.heygen_api_key = heygen_api_key or os.getenv("HEYGEN_API_KEY")
        self.openai_api_key = openai_api_key or os.getenv("OPENAI_API_KEY")
        # User can provide their own avatar and voice IDs, if not provided it will use default IDs from .env file.
        self.avatar_id = pptx_avatar_id or os.getenv("HEYGEN_AVATAR_ID")
        self.voice_id = pptx_voice_id or os.getenv("HEYGEN_VOICE_ID")
        self.model = model
        self.width = width
        self.height = height
        self.background_color = background_color
        self.request_timeout_s = request_timeout_s
        self.poll_interval_s = poll_interval_s
        # MODIFICATION: Hardcoded avatar position
        self.use_slides_as_background = use_slides_as_background
        self.slide_asset_ids: List[str] = []

        if not self.heygen_api_key:
            raise ValueError("Missing HEYGEN_API_KEY")
        if not self.openai_api_key:
            raise ValueError("Missing OPENAI_API_KEY")
        if not self.avatar_id:
            raise ValueError("Missing HEYGEN_AVATAR_ID")
        if not self.voice_id:
            raise ValueError("Missing HEYGEN_VOICE_ID")

        self._client = OpenAI(api_key=self.openai_api_key)
        self._heygen_base_v2 = "https://api.heygen.com/v2"
        self._headers = {
            "Accept": "application/json",
            "Content-Type": "application/json",
            "X-Api-Key": self.heygen_api_key,
        }

    def _post_with_retry(self, url: str, payload: Dict) -> requests.Response:
        for attempt in range(3):
            try:
                resp = requests.post(url, headers=self._headers, json=payload, timeout=self.request_timeout_s)
                resp.raise_for_status()
                return resp
            except requests.RequestException as e:
                logging.warning(f"POST attempt {attempt + 1} failed: {e}")
                if attempt == 2:
                    raise
                time.sleep(2 ** attempt)
        raise RuntimeError("Unreachable")

    def _get_with_retry(self, url: str) -> requests.Response:
        for attempt in range(3):
            try:
                resp = requests.get(url, headers=self._headers, timeout=self.request_timeout_s)
                resp.raise_for_status()
                return resp
            except requests.RequestException as e:
                logging.warning(f"GET attempt {attempt + 1} failed: {e}")
                if attempt == 2:
                    raise
                time.sleep(2 ** attempt)
        raise RuntimeError("Unreachable")

    def _create_heygen_folder(self, folder_name: str) -> str:
        logging.info(f"Creating HeyGen asset folder: {folder_name}")
        payload = {"name": folder_name}
        # The URL for the create folder endpoint is v1
        url = "https://api.heygen.com/v1/folders/create"
        response = self._post_with_retry(url, payload)
        
        # The key in the response data is 'id', not 'folder_id'.
        folder_id = response.json().get("data", {}).get("id")

        if not folder_id:
            raise RuntimeError(f"Failed to create HeyGen folder: {response.text}")
        logging.info(f"Successfully created folder with ID: {folder_id}")
        return folder_id

    def _upload_asset_to_heygen(self, file_path: str, folder_id: str) -> str:
        file_name = os.path.basename(file_path)
        mime_type, _ = mimetypes.guess_type(file_path)
        if not mime_type:
            mime_type = "image/jpeg"  # Default MIME type

        logging.info(f"Uploading {file_name} to HeyGen (target folder: {folder_id})...")

        url = "https://upload.heygen.com/v1/asset"
        
        # Create headers specific to this file upload request
        upload_headers = {
            "X-Api-Key": self.heygen_api_key,
            "Content-Type": mime_type,
        }

        with open(file_path, 'rb') as f:
            for attempt in range(3):
                try:
                    f.seek(0)
                    resp = requests.post(url, headers=upload_headers, data=f, timeout=self.request_timeout_s)
                    resp.raise_for_status()
                    
                    response_data = resp.json()
                    asset_id = response_data.get("data", {}).get("id")

                    if not asset_id:
                        raise RuntimeError(f"Failed to get asset_id from HeyGen response: {response_data}")

                    logging.info(f"Successfully uploaded {file_name}. Asset ID: {asset_id}")
                    return asset_id

                except requests.RequestException as e:
                    logging.warning(f"Upload attempt {attempt + 1} for {file_name} failed: {e}")
                    if attempt == 2:
                        raise
                    time.sleep(2 ** attempt)
        
        raise RuntimeError(f"File upload failed for {file_name} after multiple retries.")

    def _extract_slides_as_images_and_upload(self, pptx_path: str, prs: Presentation, folder_id: str):
        import tempfile
        import sys
        
        temp_dir = tempfile.mkdtemp()
        logging.info(f"Created temporary directory for slide images: {temp_dir}")
        
        slide_image_paths = []
        try:
            if sys.platform == "win32":
                import win32com.client
                logging.info("Using PowerPoint COM automation to export slides...")
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                presentation = powerpoint.Presentations.Open(os.path.abspath(pptx_path), ReadOnly=True, WithWindow=False)
                num_slides_to_process = len(prs.slides)
                for i in range(1, num_slides_to_process + 1):
                    path = os.path.join(temp_dir, f"slide_{i}.png")
                    presentation.Slides(i).Export(path, "PNG", self.width, self.height)
                    slide_image_paths.append(path)
                presentation.Close()
                powerpoint.Quit()
            else:
                raise NotImplementedError("Automatic slide export is only supported on Windows.")
        except Exception as e:
            logging.warning(f"Slide export failed: {e}. Falling back to creating blank images with text.")
            slide_image_paths = self._create_fallback_slide_images(prs, temp_dir)

        self.slide_asset_ids = [self._upload_asset_to_heygen(p, folder_id) for p in slide_image_paths]

    def _create_fallback_slide_images(self, prs: Presentation, temp_dir: str) -> List[str]:
        from PIL import Image, ImageDraw, ImageFont
        
        paths = []
        for i, slide in enumerate(prs.slides):
            path = os.path.join(temp_dir, f"slide_{i + 1}.png")
            img = Image.new("RGB", (self.width, self.height), color=self.background_color.lstrip('#'))
            draw = ImageDraw.Draw(img)
            try:
                font = ImageFont.truetype("arial.ttf", 24)
            except IOError:
                font = ImageFont.load_default()
            
            text = "\n".join(s.text for s in slide.shapes if hasattr(s, "text") and s.text).strip()
            draw.text((50, 50), f"Slide {i + 1}\n\n{text}", fill="black", font=font)
            img.save(path)
            paths.append(path)
        return paths

    def _build_video_inputs(self, slide_notes: List[str]) -> List[Dict]:
        position_map = {
            "right_corner": {"x": 0.8, "y": 0.85, "scale": 0.35},
            "left_corner": {"x": 0.2, "y": 0.85, "scale": 0.35},
            "center": {"x": 0.5, "y": 0.5, "scale": 0.8},
            "full": {"x": 0.5, "y": 0.5, "scale": 1.0},
        }


        video_inputs: List[Dict] = []
        for idx, note in enumerate(slide_notes):
            scene = {
                "character": {
                    "type": "talking_photo",
                    "talking_photo_id": self.avatar_id,
                    "talking_photo_style": "circle",
                    "scale":  0.33,
                    "offset": {
                        "x": 0.42,
                        "y": 0.42
                        },
                },
                "voice": {"type": "text", "input_text": note, "voice_id": self.voice_id},
            }

            if self.use_slides_as_background and idx < len(self.slide_asset_ids):
                asset_id = self.slide_asset_ids[idx]
                logging.info(f"Setting background for scene {idx + 1} with asset_id: {asset_id}")
                scene["background"] = {"type": "image", "image_asset_id": asset_id}
            else:
                scene["background"] = {"type": "color", "value": self.background_color}
            
            video_inputs.append(scene)
        return video_inputs

    def create_multi_scene_video(self, slide_notes: List[str], title: str) -> Dict:
        video_inputs = self._build_video_inputs(slide_notes)
        payload = {
            "video_inputs": video_inputs,
            "dimension": {"width": self.width, "height": self.height},
            "title": title,
        }
        generate_url = f"{self._heygen_base_v2}/video/generate"
        resp = self._post_with_retry(generate_url, payload)
        data = resp.json()
        video_id = data.get("data", {}).get("video_id")
        logging.info(f"Generating video using video_id : {video_id}, with Avatar id {self.avatar_id} and voice id {self.voice_id}.")
        if not video_id:
            raise RuntimeError(f"Unexpected response from generate endpoint: {resp.text}")
        return {"video_id": video_id}

    def wait_for_video(self, video_id: str) -> Dict:        
        """
        Polls the HeyGen API for the video status indefinitely until it is
        either 'completed' or 'failed'.
        """
        logging.info(f"Waiting for video {video_id} to complete. with Avatar id {self.avatar_id} and voice id {self.voice_id}. Polling status...")

        # Loop indefinitely until a terminal status is reached
        while True:
            status_url = f"https://api.heygen.com/v1/video_status.get?video_id={video_id}"
            try:
                resp = self._get_with_retry(status_url)
                data = resp.json().get("data", {})
                status = data.get("status")

                if status in ("completed", "success"):
                    logging.info("Video generation completed successfully.")
                    return {"status": "completed", "video_url": data.get("video_url")}
                
                elif status in ("failed", "error"):
                    logging.error(f"Video generation failed with data: {data}")
                    raise RuntimeError(f"Video generation failed: {data}")
                
                else:
                    logging.info(f"Video status is '{status}', continuing to wait...")

            except requests.RequestException as e:
                logging.warning(f"Status check failed due to a network error: {e}. Retrying...")
            
            # Wait for the specified interval before the next poll
            time.sleep(self.poll_interval_s)
    
    def generate_speaker_notes(self, slide_texts: List[str]) -> List[str]:
        logging.info(f"Generating speaker notes for {len(slide_texts)} slides using model {self.model}...")
        
        speaker_notes = []
        for i, text in enumerate(slide_texts):
            try:
                system_prompt = (
                    "You are a virtual teacher, and your task is to explain the content of a presentation slide. "
                    "You will be given the text from a slide and are expected to generate a script that is clear, "
                    "engaging, and educational. Your tone should be professional yet approachable."
                )
                
                user_prompt = f"Here is the slide content:\n\n---\n\n{text}\n\n---\n\nPlease provide the speaker notes."
                
                response = self._client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt},
                    ],
                    temperature=0.7,
                    max_tokens=250,
                )
                
                note = response.choices[0].message.content.strip()
                speaker_notes.append(note)
                logging.info(f"Generated note for slide {i+1}: '{note[:100]}...'")

            except Exception as e:
                logging.error(f"Failed to generate notes for slide {i+1}: {e}")
                # Fallback to a simple note
                speaker_notes.append(f"This is the content for slide {i+1}: {text}")

        return speaker_notes

    def convert(self, pptx_path: str, title: Optional[str] = None, max_slides: Optional[int] = None) -> Dict:
        logging.info(f"Starting conversion for: {pptx_path}")
        
        prs = Presentation(pptx_path)
        slides = list(prs.slides)
        
        if not slides:
            raise ValueError("No slides found in the PowerPoint file.")

        if max_slides and len(slides) > max_slides:
            logging.warning(f"Limiting to {max_slides} slides from {len(slides)}.")
            slides = slides[:max_slides]
        
        slides_text = ["\n".join(s.text for s in slide.shapes if hasattr(s, "text") and s.text).strip() for slide in slides]

        if self.use_slides_as_background:
            folder_name = f"Slides_{os.path.basename(pptx_path)}_{int(time.time())}"
            folder_id = self._create_heygen_folder(folder_name)
            self._extract_slides_as_images_and_upload(pptx_path, prs, folder_id)
            if len(self.slide_asset_ids) < len(slides_text):
                logging.warning("Number of uploaded slide assets is less than the number of slides.")
                slides_text = slides_text[:len(self.slide_asset_ids)]

        slide_notes = self.generate_speaker_notes(slides_text)
        
        final_title = title or os.path.basename(pptx_path)
        gen_info = self.create_multi_scene_video(slide_notes, title=final_title)
        final_status = self.wait_for_video(gen_info["video_id"])

        return {
            "video_id": gen_info["video_id"],
            "video_url": final_status.get("video_url"),
            "slides_count": len(slide_notes),
        }

def main():
    pptx_file_path = input("enter ppt file path here :")
    pptx_avatar_id = input("enter ppt avatar id here :")
    pptx_voice_id = input("enter ppt voice id here :")

    if not os.path.exists(pptx_file_path):
        logging.error(f"File not found: {pptx_file_path}")
        logging.error("Please update the 'pptx_file_path' variable in the main() function with the correct path.")
        return

    try:
        # --- MODIFICATION: Simplified converter initialization ---
        # Avatar ID and Voice ID are now handled by environment variables within the class.
        # Using slides as background is now the default.
        converter = PPTXToHeyGenVideo()
        
        # --- MODIFICATION: Simplified convert method call ---
        # Title and max_slides are no longer passed as command-line arguments.
        result = converter.convert(pptx_file_path)
        
        print("\n--- Conversion Successful ---")
        print(json.dumps(result, indent=2))
    except Exception as e:
        logging.error(f"An error occurred during conversion: {e}", exc_info=True)

if __name__ == "__main__":
    main()